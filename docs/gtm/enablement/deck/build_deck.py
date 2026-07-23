"""Build the 'Directing Claude' training deck from deck.md into an editable .pptx.

Usage:
    python build_deck.py                 # deck.md -> Directing-Claude.pptx (this folder)
    python build_deck.py IN OUT [LOGO]   # custom paths

deck.md format (slides separated by a line containing only '---'):

    # Slide title
    @layout: title | section | content        (optional; default content)
    - a bullet
    - another bullet
    @notes:
    Speaker notes, one or more lines, to end of slide.
"""
from __future__ import annotations

import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).resolve().parent
DEFAULT_IN = HERE / "deck.md"
DEFAULT_OUT = HERE / "Directing-Claude.pptx"
DEFAULT_LOGO = HERE / "assets" / "harbormill-logo.png"

# Harbormill brand: black wordmark on a warm light plate. Light reads better projected.
INK = RGBColor(0x0A, 0x0A, 0x0A)
MUTED = RGBColor(0x55, 0x55, 0x55)
BG = RGBColor(0xF7, 0xF6, 0xF2)

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


@dataclass
class Slide:
    title: str = ""
    layout: str = "content"  # title | section | content
    bullets: list[str] = field(default_factory=list)
    notes: str = ""


def parse_deck(md: str) -> list[Slide]:
    slides: list[Slide] = []
    for block in re.split(r"(?m)^---\s*$", md):
        if not block.strip():
            continue
        slide = Slide()
        note_lines: list[str] = []
        in_notes = False
        for line in block.splitlines():
            if line.strip() == "@notes:":
                in_notes = True
                continue
            if in_notes:
                note_lines.append(line)
                continue
            layout_m = re.match(r"@layout:\s*(\w+)", line.strip())
            if layout_m:
                slide.layout = layout_m.group(1).lower()
                continue
            if line.startswith("# "):
                slide.title = line[2:].strip()
                continue
            bullet_m = re.match(r"\s*[-*]\s+(.*)", line)
            if bullet_m:
                slide.bullets.append(bullet_m.group(1).strip())
                continue
            if line.strip():
                slide.bullets.append(line.strip())
        slide.notes = "\n".join(note_lines).strip()
        slides.append(slide)
    return slides


def _fill_background(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def _add_title(slide, text: str, *, size: int, top: float, align=PP_ALIGN.LEFT,
               left: float = 0.7, width: float = 12.0) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = INK


def _add_bullets(slide, bullets: list[str], *, top: float = 2.3) -> None:
    if not bullets:
        return
    box = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.5), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "•  " + b
        run.font.size = Pt(22)
        run.font.color.rgb = INK


def _add_logo(slide, logo_path: Path, *, kind: str) -> None:
    if not Path(logo_path).exists():
        return
    if kind == "title":
        slide.shapes.add_picture(str(logo_path), Inches(4.67), Inches(1.3), height=Inches(2.0))
    elif kind == "section":
        slide.shapes.add_picture(str(logo_path), Inches(0.6), Inches(0.5), height=Inches(0.8))
    else:  # content — small, top-right
        slide.shapes.add_picture(str(logo_path), Inches(11.6), Inches(0.4), height=Inches(0.6))


def build(md_path, out_path, logo_path=DEFAULT_LOGO) -> None:
    md = Path(md_path).read_text(encoding="utf-8")
    slides = parse_deck(md)

    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]

    for s in slides:
        slide = prs.slides.add_slide(blank)
        _fill_background(slide)
        if s.layout == "title":
            _add_logo(slide, logo_path, kind="title")
            _add_title(slide, s.title, size=44, top=3.6, align=PP_ALIGN.CENTER, left=0.7, width=12.0)
            if s.bullets:
                _add_title(slide, s.bullets[0], size=20, top=4.9, align=PP_ALIGN.CENTER)
        elif s.layout == "section":
            _add_logo(slide, logo_path, kind="section")
            _add_title(slide, s.title, size=40, top=3.2, align=PP_ALIGN.CENTER, left=0.7, width=12.0)
        else:
            _add_logo(slide, logo_path, kind="content")
            _add_title(slide, s.title, size=32, top=0.5)
            _add_bullets(slide, s.bullets, top=2.0)
        if s.notes:
            slide.notes_slide.notes_text_frame.text = s.notes

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def main(argv: list[str]) -> None:
    in_path = Path(argv[1]) if len(argv) > 1 else DEFAULT_IN
    out_path = Path(argv[2]) if len(argv) > 2 else DEFAULT_OUT
    logo_path = Path(argv[3]) if len(argv) > 3 else DEFAULT_LOGO
    build(in_path, out_path, logo_path)
    print(f"Built {out_path} from {in_path}")


if __name__ == "__main__":
    main(sys.argv)
