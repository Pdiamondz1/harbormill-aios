from pathlib import Path

from pptx import Presentation

from build_deck import parse_deck, build

HERE = Path(__file__).resolve().parent
LOGO = HERE / "assets" / "harbormill-logo.png"

SAMPLE = """\
# Directing Claude
@layout: title
- A training deck
@notes:
Welcome the room. State the one idea.
---
# Session 1
@layout: section
---
# Get the tools
- Install Claude Code
- Sign in
@notes:
Walk them through sign-in slowly.
"""


def test_parse_deck_splits_and_reads_fields():
    slides = parse_deck(SAMPLE)
    assert len(slides) == 3
    assert slides[0].title == "Directing Claude"
    assert slides[0].layout == "title"
    assert slides[1].layout == "section"
    assert slides[2].bullets == ["Install Claude Code", "Sign in"]
    assert "sign-in slowly" in slides[2].notes


def test_build_writes_pptx_with_expected_slide_count(tmp_path):
    md = tmp_path / "deck.md"
    md.write_text(SAMPLE, encoding="utf-8")
    out = tmp_path / "out.pptx"
    build(md, out, LOGO)
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == 3


def test_title_slide_has_logo_picture(tmp_path):
    md = tmp_path / "deck.md"
    md.write_text(SAMPLE, encoding="utf-8")
    out = tmp_path / "out.pptx"
    build(md, out, LOGO)
    prs = Presentation(str(out))
    # picture shape type == 13 (MSO_SHAPE_TYPE.PICTURE)
    pics = [sh for sh in prs.slides[0].shapes if sh.shape_type == 13]
    assert len(pics) >= 1


def test_notes_are_written(tmp_path):
    md = tmp_path / "deck.md"
    md.write_text(SAMPLE, encoding="utf-8")
    out = tmp_path / "out.pptx"
    build(md, out, LOGO)
    prs = Presentation(str(out))
    assert "Welcome the room" in prs.slides[0].notes_slide.notes_text_frame.text
